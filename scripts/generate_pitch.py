
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION

# --- Configuration ---
FILENAME = "Diablex_iHub_S4_Pitch_Deck.pptx"

# Colors
BG_COLOR = RGBColor(20, 30, 40)       # Dark Navy / Charcoal
ACCENT_COLOR = RGBColor(0, 255, 200)  # Electric Teal
TEXT_WHITE = RGBColor(255, 255, 255)
TEXT_GRAY = RGBColor(200, 200, 200)

def set_slide_background(slide):
    """Sets the slide background to the dark theme."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_title(slide, text, font_size=40):
    """Adds a styled title to the slide."""
    # Create the title textbox manually
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.5))
    title_shape.text_frame.text = text
    
    tfr = title_shape.text_frame
    p = tfr.paragraphs[0]
    p.font.color.rgb = ACCENT_COLOR
    p.font.name = "Arial"
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

def add_content(slide, text_lines, top=2.0, font_size=24):
    """Adds bulleted content to the slide."""
    left = Inches(0.5)
    width = Inches(9)
    height = Inches(5)
    
    txBox = slide.shapes.add_textbox(left, Inches(top), width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(text_lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = line
        p.font.color.rgb = TEXT_WHITE
        p.font.name = "Arial"
        p.font.size = Pt(font_size)
        p.space_after = Pt(14)

def add_speaker_notes(slide, notes):
    """Adds speaker notes to the slide."""
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes

# --- Slide Generators ---

def create_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    set_slide_background(slide)
    
    # Title
    title = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Diablex"
    p.font.color.rgb = ACCENT_COLOR
    p.font.size = Pt(80)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    tf_sub = sub.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "The Digital Co-Pilot for Predictive Diabetes Care"
    p_sub.font.color.rgb = TEXT_WHITE
    p_sub.font.size = Pt(32)
    p_sub.alignment = PP_ALIGN.CENTER

    add_speaker_notes(slide, "Good morning. I am here to introduce Diablex. We are building the future of predictive diabetes care. We are not just another logging app; we are a digital co-pilot designed to anticipate metabolic changes before they happen.")

def create_problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "The Hook: Data Fatigue")
    
    content = [
        "• Patients generate 288+ glucose points daily but lack insight.",
        "• Doctors have <15 mins to review months of logs.",
        "• Result: Reactive care, anxiety, and 'Burnout'.",
        "• The Gap: Abundant data, zero foresight."
    ]
    add_content(slide, content)
    
    # Visual abstraction (Example: "Cloud" of numbers obscuring a person)
    # Simple representation: A messy pile of shapes vs clear arrow
    
    add_speaker_notes(slide, "The problem isn't a lack of data; it's data fatigue. Patients and doctors are drowning in logs but starving for insights. This leads to reactive care—fixing issues after they occur—rather than preventing them.")

def create_solution_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "The Solution: 'Antigravity'")
    
    content = [
        "• Moving from Reactive to Predictive.",
        "• Automates the cognitive load of diabetes management.",
        "• Live alerts for UPCOMING glucose events.",
        "• Seamless integration of Hardware + AI."
    ]
    add_content(slide, content)
    
    add_speaker_notes(slide, "Our solution is 'Antigravity.' We lift the mental burden. By closing the loop between IoT collection and AI analysis, we provide real-time, predictive alerts. It's not just about what is happening now, but what will happen in 30 minutes.")

def create_stack_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "How it Works: The Stack")
    
    # Diagram: Hardware -> Backend -> Intelligence -> User
    
    y = Inches(3.5)
    box_w = Inches(1.8)
    box_h = Inches(1)
    gap = Inches(0.5)
    start_x = Inches(0.5)
    
    steps = ["Hardware (IoT)", "Backend (Node.js)", "Intelligence (Py/Pandas)", "User (React/Vite)"]
    
    for i, step in enumerate(steps):
        x = start_x + (i * (box_w + gap))
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_COLOR if i == 2 else RGBColor(40, 50, 60) # Highlight Intelligence
        shape.shadow.inherit = False
        
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = step
        p.font.color.rgb = TEXT_WHITE if i != 2 else RGBColor(0,0,0)
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER
        
        # Arrow
        if i < len(steps) - 1:
            arrow_x = x + box_w
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, y + (box_h/2) - Inches(0.1), gap, Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = TEXT_GRAY

    add_speaker_notes(slide, "Here is our technical architecture. A custom IoT sensor feeds data to a high-speed Node.js backend. The core intelligence resides in our Python/Pandas engine, which runs our predictive models. The insights are delivered instantly to the user via a modern React application.")

def create_engine_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "The Secret Sauce: Predictive Intelligence")
    
    content = [
        "• ARIMA & Time-Series Forecasting Models.",
        "• Predicts hypoglycemia 20-30 mins in advance.",
        "• Personalized algorithm adapts to user metabolism.",
        "• Foresight > Hindsight."
    ]
    add_content(slide, content, top=4.5)
    
    # Chart Simulation
    chart_data = CategoryChartData()
    chart_data.categories = ['10:00', '10:10', '10:20', '10:30', '10:40', '10:50']
    series_actual = chart_data.add_series('Actual', (120, 115, 110, None, None, None))
    series_pred = chart_data.add_series('Diablex Prediction', (None, None, 110, 95, 80, 70))
    
    x, y, w, h = Inches(0.5), Inches(1.5), Inches(6), Inches(3)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, w, h, chart_data).chart
    
    # Style chart (Basic styling availability depends on pptx version, keeping it simple)
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    add_speaker_notes(slide, "Our core differentiator is the Predictive Engine. Using ARIMA and specific time-series models, we can project glucose trajectories. In this graph, you see the solid line as current data, and the dotted line as our prediction—alerting the user of a drop 20 minutes before it actually happens.")

def create_market_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "Market Strategy")
    
    content = [
        "• Primary: Type 1 Diabetics (Insulin Dependent).",
        "• Secondary: Endocrinologists & Clinics.",
        "• Focus: Gujarat Healthcare Ecosystem (scalable model).",
        "• Growth: B2C App Subscription + B2B Clinic Dashboards."
    ]
    add_content(slide, content)
    
    add_speaker_notes(slide, "We are targeting the high-need segment: Type 1 diabetics and insulin-dependent Type 2s. Our go-to-market starts here in the Gujarat ecosystem, leveraging local clinical partnerships before scaling nationally.")

def create_ask_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "The Ask & Roadmap")
    
    content = [
        "• Grant Request: 2.5 – 10 Lakh INR.",
        "• Usage:",
        "    - Prototype Refinement (HW Miniaturization).",
        "    - Cloud Infrastructure Scaling.",
        "    - Initial Clinical Data Validation.",
        "• Goal: Market-ready MVP within 6 months."
    ]
    add_content(slide, content)
    
    add_speaker_notes(slide, "We are seeking seed support to move from a working prototype to a field-ready device. The funds will drive miniaturization and essential validation studies. Thank you.")

# --- Main ---

def main():
    prs = Presentation()
    
    # Slide 1: Title
    create_title_slide(prs)
    
    # Slide 2: The Hook
    create_problem_slide(prs)
    
    # Slide 3: The Solution
    create_solution_slide(prs)
    
    # Slide 4: The Stack
    create_stack_slide(prs)
    
    # Slide 5: The Secret Sauce
    create_engine_slide(prs)
    
    # Slide 6: Market Strategy
    create_market_slide(prs)
    
    # Slide 7: Ask
    create_ask_slide(prs)
    
    prs.save(FILENAME)
    print(f"Presentation saved to {FILENAME}")

if __name__ == "__main__":
    main()
